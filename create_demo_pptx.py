from pptx import Presentation
from pptx.util import Inches, Pt

# Paths
template_path = r"C:\Users\JOBLER\OneDrive - Sleep Number Corporation\Documents\AI TPRM MACHINE\Jan2026_AIHackathon_PresentationTemplate.pptx"
output_path = r"C:\Users\JOBLER\OneDrive - Sleep Number Corporation\Documents\AI TPRM MACHINE\TPRM_Demo_Presentation.pptx"

# Load template
prs = Presentation(template_path)

# Helper function to add a slide with title and bullet content
def add_content_slide(prs, title, bullets, layout_index=1):
    slide_layout = prs.slide_layouts[layout_index]
    slide = prs.slides.add_slide(slide_layout)
    if slide.shapes.title:
        slide.shapes.title.text = title
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 1:
            tf = shape.text_frame
            tf.clear()
            for i, bullet in enumerate(bullets):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = bullet
                p.level = 0
            break
    return slide

# Clear existing slides
while len(prs.slides) > 0:
    rId = prs.slides._sldIdLst[0].rId
    prs.part.drop_rel(rId)
    del prs.slides._sldIdLst[0]

# Slide 1: Title
slide = prs.slides.add_slide(prs.slide_layouts[0])
if slide.shapes.title:
    slide.shapes.title.text = "AI-Powered Third Party Risk Management"
for shape in slide.placeholders:
    if shape.placeholder_format.idx == 1:
        shape.text = "4-Minute Demo"
        break

# Slide 2: Demo Flow
add_content_slide(prs, "Demo: What You'll See", [
    "Dashboard - Real-time risk posture at a glance",
    "AI Vendor Profiling - Instant risk classification",
    "Document Analysis - Automated finding extraction",
    "Reporting - Executive-ready insights in seconds"
])

# Slide 3: Six AI Agents
add_content_slide(prs, "How AI Helps: Six Specialized Agents", [
    "Vendor Profiler - Auto-classifies risk level instantly",
    "Security Analyzer - Extracts findings from documents",
    "Risk Assessor - Evaluates 9 security domains",
    "Risk Manager - Generates remediation plans",
    "Report Generator - Creates executive summaries",
    "Information Gatherer - Identifies documentation gaps"
])

# Slide 4: Business Value
add_content_slide(prs, "Business Value", [
    "90% reduction in manual document review time",
    "Vendor onboarding reduced from days to hours",
    "Consistent risk scoring across all vendors",
    "Auto-mapping to SOC 2, ISO 27001, NIST, HIPAA, PCI-DSS",
    "Scalable - more vendors without adding headcount"
])

# Slide 5: Summary
add_content_slide(prs, "Summary", [
    "Transforms vendor risk from reactive to proactive",
    "AI handles review, assessment, and reporting",
    "Team focuses on decisions, not document review",
    "Reduces risk | Ensures compliance | Improves efficiency"
])

# Save
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
